import re
from langchain_community.document_loaders import PyPDFLoader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import requests
from bs4 import BeautifulSoup
from readability import Document as ReadabilityDocument
import io
import os
from dotenv import load_dotenv

# --- PDF Text Extraction ---
def extract_text_from_pdf(uploaded_file):
    with open("temp_uploaded_file.pdf", "wb") as f:
        f.write(uploaded_file.getbuffer())

    pdf_loader = PyPDFLoader("temp_uploaded_file.pdf")
    docs = pdf_loader.load()

    unwanted_patterns = [
        r'.*indd.*',
        r'^\s+$'
    ]

    full_text = ""
    for page_number, doc in enumerate(docs, start=1):
        page_content = doc.page_content
        filtered_content = []

        for line in page_content.split('\n'):
            if any(re.match(pattern, line.strip()) for pattern in unwanted_patterns):
                continue
            filtered_content.append(line.strip())

        for line in filtered_content:
            if line.strip():
                full_text += line + "\n"

    return full_text

# --- TXT File Text Extraction ---
def extract_text_from_txt(uploaded_file):
    lines = uploaded_file.read().decode("utf-8").splitlines()

    unwanted_patterns = [
        r'.*indd.*',
        r'^\s+$'
    ]

    full_text = ""
    for line in lines:
        stripped_line = line.strip()
        if any(re.match(pattern, stripped_line) for pattern in unwanted_patterns):
            continue
        if stripped_line:
            full_text += stripped_line + "\n"

    return full_text

# --- DOCX File Text Extraction ---
def extract_text_from_docx(uploaded_file):
    doc = Document(uploaded_file)

    unwanted_patterns = [
        r'.*indd.*',
        r'^\s+$'
    ]

    full_text = ""
    for para in doc.paragraphs:
        line = para.text.strip()
        if any(re.match(pattern, line) for pattern in unwanted_patterns):
            continue
        if line:
            full_text += line + "\n"

    return full_text

# --- PPTX File Text Extraction ---
def extract_text_from_pptx(uploaded_file):
    prs = Presentation(uploaded_file)

    unwanted_patterns = [
        r'.*indd.*',
        r'^\s+$'
    ]

    full_text = ""
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for line in shape.text.split('\n'):
                    stripped_line = line.strip()
                    if any(re.match(pattern, stripped_line) for pattern in unwanted_patterns):
                        continue
                    if stripped_line:
                        full_text += stripped_line + "\n"

    return full_text

# --- Image Text Extraction ---
def extract_text_from_image(uploaded_file):
    api_key = os.getenv("OCR_API_KEY")  
    
    if not api_key:
        return "Error: OCR API key is missing."

    url = 'https://api.ocr.space/parse/image'
    
    try:
        image = Image.open(uploaded_file)
    except Exception as e:
        return f"Error: Unable to open image file. {str(e)}"
    
    img_byte_arr = io.BytesIO()
    image.save(img_byte_arr, format='PNG')  # Convert to PNG format
    img_byte_arr = img_byte_arr.getvalue()  # Get the byte data

    try:
        response = requests.post(
            url,
            files={'filename': ('image.png', img_byte_arr, 'image/png')},
            data={'apikey': api_key, 'language': 'eng'}
        )
    except requests.RequestException as e:
        return f"Error: Request to OCR API failed. {str(e)}"
    
    try:
        result = response.json()
    except ValueError:
        return f"Error: Response from OCR.Space is not in valid JSON format. Response: {response.text}"

    if result.get("IsErroredOnProcessing"):
        return f"OCR failed: {result.get('ErrorMessage', ['Unknown error'])[0]}"
    
    extracted_text = result["ParsedResults"][0]["ParsedText"]
    
    unwanted_patterns = [
        r'.*indd.*',   
        r'^\s+$'       
    ]
    
    full_text = ""
    for line in extracted_text.split('\n'):
        if any(re.match(pattern, line.strip()) for pattern in unwanted_patterns):
            continue
        if line.strip():  
            full_text += line.strip() + "\n"
    
    return full_text

# --- URL Article Extraction ---
def extract_text_from_url(url):
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
            "(KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36"
        )
    }

    try:
        response = requests.get(url, headers=headers)
        response.raise_for_status()

        doc = ReadabilityDocument(response.text)
        title = doc.title()
        summary_html = doc.summary()

        soup = BeautifulSoup(summary_html, 'html.parser')
        text = soup.get_text(separator='\n')

        unwanted_patterns = [
            r'.*indd.*',
            r'^\s*$'
        ]

        full_text = f"Title: {title}\n\n"

        for line in text.split('\n'):
            if not any(re.match(p, line.strip()) for p in unwanted_patterns):
                full_text += line.strip() + "\n"

        return full_text

    except Exception as e:
        return f"Error extracting article: {e}"
